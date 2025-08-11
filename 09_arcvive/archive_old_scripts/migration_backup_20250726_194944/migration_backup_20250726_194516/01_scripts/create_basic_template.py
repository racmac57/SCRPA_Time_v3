# 🕒 2025-01-28-00-02-00
# SCRPA_Fix/create_basic_template.py
# Author: R. A. Carucci
# Purpose: Create a basic PowerPoint template for SCRPA reports

from pptx import Presentation
from pptx.util import Inches
import os

def create_basic_scrpa_template():
    """Create a basic SCRPA PowerPoint template"""
    
    # Create new presentation
    prs = Presentation()
    
    # Define crime types
    crime_types = [
        "MV Theft",
        "Burglary Auto", 
        "Burglary - Comm & Res",
        "Robbery",
        "Sexual Offenses"
    ]
    
    # Create slides for each crime type
    for i, crime_type in enumerate(crime_types):
        # Add slide
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(12), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"SCRPA 7-Day Analysis - {crime_type}"
        
        # Add placeholder for chart (left side)
        chart_placeholder = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(5.84), Inches(4)
        )
        chart_frame = chart_placeholder.text_frame
        chart_frame.text = f"[{crime_type} Chart will be inserted here]"
        
        # Add placeholder for map (right side)
        map_placeholder = slide.shapes.add_textbox(
            Inches(6.8), Inches(2), Inches(5.6), Inches(4.03)
        )
        map_frame = map_placeholder.text_frame  
        map_frame.text = f"[{crime_type} Map will be inserted here]"
        
        # Add placeholder for incident table (top right)
        table_placeholder = slide.shapes.add_textbox(
            Inches(6.8), Inches(0.8), Inches(8.33), Inches(1.98)
        )
        table_frame = table_placeholder.text_frame
        table_frame.text = f"[{crime_type} Incident Table will be inserted here]"
    
    return prs

# Create and save template
if __name__ == "__main__":
    print("Creating basic SCRPA PowerPoint template...")
    
    # Create template
    template = create_basic_scrpa_template()
    
    # Save to expected location
    template_dir = r"C:\Users\carucci_r\OneDrive - City of Hackensack\_25_SCRPA\TIME_Based\7_Day_Templet_SCRPA_Time"
    template_path = os.path.join(template_dir, "SCRPA_REPORT.pptx")
    
    # Ensure directory exists
    os.makedirs(template_dir, exist_ok=True)
    
    # Save template
    template.save(template_path)
    
    print(f"✅ Template created: {template_path}")
    print("✅ You can now run your SCRPA script!")
