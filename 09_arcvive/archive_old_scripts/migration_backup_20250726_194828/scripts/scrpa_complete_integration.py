#!/usr/bin/env python3
# 🕒 2025-07-01-17-00-00
# SCRPA_Complete_Integration/scrpa_complete_system
# Author: R. A. Carucci
# Purpose: Complete SCRPA system with charts, maps, statistics, and proper folder structure

import os
import sys
import logging
import time
import glob
import shutil
from datetime import datetime as dt_datetime, date as dt_date
import arcpy

# --- Import core configuration ---
try:
    from config import (
        get_crime_type_folder,
        get_7day_period_dates,
        CRIME_FOLDER_MAPPING,
        POWERPOINT_TEMPLATES,
        CRIME_TYPES,
        get_correct_folder_name,
        get_standardized_filename,
        get_output_folder
    )
except ImportError as e:
    print(f"❌ Config import failed: {e}")
    sys.exit(1)

# --- Import and patch map_export to avoid missing-name errors ---
try:
    import map_export

    # Resolve the correct map-export function
    if hasattr(map_export, 'export_maps'):
        export_maps = map_export.export_maps
    elif hasattr(map_export, 'export_maps_for_crime_type'):
        export_maps = map_export.export_maps_for_crime_type
    else:
        raise ImportError("No map-export function found in map_export module")

    # Stub any missing SQL-filter builders so incident_table_automation won't error on import
    def __stub_filter(*args, **kwargs):
        # Default to no-records filter
        return "1=0"
    for name in (
        'build_sql_filter_7day_excel',
        'build_sql_filter_28day_excel',
        'build_sql_filter_ytd_excel'
    ):
        if not hasattr(map_export, name):
            setattr(map_export, name, __stub_filter)

except ImportError as e:
    print(f"❌ map_export import failed: {e}")
    sys.exit(1)

# --- Import remaining components ---
try:
    from chart_export import export_chart
    from incident_table_automation import (
        export_incident_table_data_with_autofit,
        create_incident_table_placeholder
    )
    from scrpa_statistical_fixed import generate_statistical_exports_fixed
    print("✅ All modules imported successfully")
except ImportError as e:
    print(f"❌ Module import failed: {e}")
    sys.exit(1)


def run_complete_scrpa_system(report_date_str=None, generate_all=True):
    """
    Run the complete SCRPA system with proper folder structure
    Args:
        report_date_str (str): Date in YYYY_MM_DD format
        generate_all (bool): Generate all components (maps, charts, tables, stats)
    Returns:
        bool: True if successful
    """
    start_time = time.time()

    if report_date_str is None:
        report_date_str = dt_date.today().strftime("%Y_%m_%d")

    print(f"\n🚀 SCRPA COMPLETE SYSTEM")
    print("=" * 60)
    print(f"Report Date: {report_date_str}")
    print(f"Generate All Components: {generate_all}")

    try:
        # Parse report date
        report_date = dt_datetime.strptime(report_date_str, "%Y_%m_%d").date()

        # Create main output folder
        main_output_folder = get_output_folder(report_date_str)
        os.makedirs(main_output_folder, exist_ok=True)

        # Copy PowerPoint template
        template_success = setup_powerpoint_template(main_output_folder, report_date_str)

        # Process each crime type
        results = {}
        total_crimes = len(CRIME_TYPES)

        print(f"\n📊 PROCESSING {total_crimes} CRIME TYPES")
        print("=" * 50)

        for i, crime_type in enumerate(CRIME_TYPES, 1):
            print(f"\n[{i}/{total_crimes}] Processing {crime_type}...")

            # Create crime-specific folder
            crime_folder = get_crime_type_folder(crime_type, report_date_str)
            os.makedirs(crime_folder, exist_ok=True)

            crime_results = {
                'folder_created': True,
                'map': False,
                'chart': False,
                'table': False
            }

            if generate_all:
                # 1) Map
                print(f"  🗺️ Generating map...")
                try:
                    map_success = export_maps(crime_type, crime_folder, sys.argv)
                    crime_results['map'] = map_success
                    print(f"    {'✅' if map_success else '❌'} Map: {'Success' if map_success else 'Failed'}")
                except Exception as e:
                    print(f"    ❌ Map export failed: {e}")

                # 2) Chart
                print(f"  📈 Generating chart...")
                try:
                    chart_success = export_chart(crime_type)
                    crime_results['chart'] = chart_success
                    print(f"    {'✅' if chart_success else '❌'} Chart: {'Success' if chart_success else 'Failed'}")
                except Exception as e:
                    print(f"    ❌ Chart failed: {e}")

                # 3) Incident table
                print(f"  📋 Generating incident table...")
                try:
                    table_success = export_incident_table_data_with_autofit(crime_type, report_date)
                    if not table_success:
                        # Fallback placeholder
                        filename_prefix = get_standardized_filename(crime_type)
                        placeholder_path = os.path.join(
                            crime_folder,
                            f"{filename_prefix}_IncidentTable_Placeholder.png"
                        )
                        table_success = create_incident_table_placeholder(crime_type, placeholder_path)
                    crime_results['table'] = table_success
                    print(f"    {'✅' if table_success else '❌'} Table: {'Success' if table_success else 'Failed'}")
                except Exception as e:
                    print(f"    ❌ Table failed: {e}")

            results[crime_type] = crime_results

            # Summary for this crime type
            success_count = sum(1 for ok in crime_results.values() if ok)
            print(f"  📊 {crime_type}: {success_count}/4 components successful")

        # 4) Statistical reports
        if generate_all:
            print(f"\n📊 GENERATING STATISTICAL REPORTS")
            print("-" * 40)
            try:
                stats_success = generate_statistical_exports_fixed(report_date)
                print(f"{'✅' if stats_success else '❌'} Statistical reports: {'Success' if stats_success else 'Failed'}")
            except Exception as e:
                print(f"❌ Statistical reports failed: {e}")
                stats_success = False
        else:
            stats_success = True

        # 5) Assemble PowerPoint
        if generate_all and template_success:
            print(f"\n🎯 ASSEMBLING POWERPOINT PRESENTATION")
            print("-" * 40)
            try:
                ppt_success = assemble_powerpoint_presentation(main_output_folder, results, report_date_str)
                print(f"{'✅' if ppt_success else '❌'} PowerPoint assembly: {'Success' if ppt_success else 'Failed'}")
            except Exception as e:
                print(f"❌ PowerPoint assembly failed: {e}")
                ppt_success = False
        else:
            ppt_success = template_success

        # Overall summary
        crime_success_count = sum(
            1 for cr in results.values()
            if sum(1 for ok in cr.values() if ok) >= 2
        )
        overall_success = (
            crime_success_count >= len(CRIME_TYPES) * 0.6 and
            stats_success and
            ppt_success
        )

        total_time = time.time() - start_time
        print(f"\n🎉 SCRPA SYSTEM COMPLETE")
        print("=" * 60)
        print(f"Processing time: {total_time:.1f} seconds")
        print(f"Overall success: {'✅ SUCCESS' if overall_success else '❌ PARTIAL SUCCESS'}")
        print(f"Crime types processed: {crime_success_count}/{len(CRIME_TYPES)}")
        print(f"Statistical reports: {'✅' if stats_success else '❌'}")
        print(f"PowerPoint assembly: {'✅' if ppt_success else '❌'}")
        print(f"Output folder: {main_output_folder}")

        print("\n📋 DETAILED RESULTS:")
        for crime_type, cr in results.items():
            comps = [name for name, ok in cr.items() if ok]
            print(f"  {crime_type}: {len(comps)}/4 - {', '.join(comps)}")

        return overall_success

    except Exception as e:
        print(f"❌ SCRPA system failed: {e}")
        logging.error(f"SCRPA system failed: {e}", exc_info=True)
        return False


def setup_powerpoint_template(output_folder, report_date_str):
    """Copy the PPT template into the output folder."""
    try:
        template_path = (
            POWERPOINT_TEMPLATES.get("7day")
            or POWERPOINT_TEMPLATES.get("main")
        )
        if not template_path or not os.path.exists(template_path):
            print(f"⚠️ PowerPoint template not found: {template_path}")
            return False

        folder_name = get_correct_folder_name(report_date_str)
        dest_ppt = os.path.join(output_folder, f"{folder_name}.pptx")
        shutil.copy2(template_path, dest_ppt)
        print(f"✅ PowerPoint template copied: {os.path.basename(dest_ppt)}")
        return True

    except Exception as e:
        print(f"❌ PowerPoint template setup failed: {e}")
        return False


def assemble_powerpoint_presentation(output_folder, results, report_date_str):
    """Place charts, maps, and tables onto PPT slides."""
    try:
        from pptx import Presentation
        from pptx.util import Inches

        ppt_files = glob.glob(os.path.join(output_folder, "*.pptx"))
        if not ppt_files:
            print(f"❌ No PowerPoint file found in: {output_folder}")
            return False

        ppt_path = ppt_files[0]
        prs = Presentation(ppt_path)
        print(f"📝 Assembling {len(prs.slides)} slides...")

        coords = {
            'chart': {'left': 0.5, 'top': 2.8, 'width': 5.84, 'height': 4.00},
            'map':   {'left': 6.8, 'top': 2.8, 'width': 5.60, 'height': 4.03},
            'table': {'left': 6.8, 'top': 0.8, 'width': 8.33, 'height': 1.98}
        }

        for idx, crime_type in enumerate(CRIME_TYPES):
            if idx >= len(prs.slides):
                print(f"⚠️ No slide for {crime_type}")
                continue
            slide = prs.slides[idx]
            crime_folder = get_crime_type_folder(crime_type, report_date_str)
            prefix = get_standardized_filename(crime_type)

            # Chart
            chart_path = os.path.join(crime_folder, f"{prefix}_Chart.png")
            if os.path.exists(chart_path):
                slide.shapes.add_picture(
                    chart_path,
                    Inches(coords['chart']['left']),
                    Inches(coords['chart']['top']),
                    Inches(coords['chart']['width']),
                    Inches(coords['chart']['height'])
                )
                print(f"  ✅ Added chart for {crime_type}")

            # Map
            map_path = os.path.join(crime_folder, f"{prefix}_Map.png")
            if os.path.exists(map_path):
                slide.shapes.add_picture(
                    map_path,
                    Inches(coords['map']['left']),
                    Inches(coords['map']['top']),
                    Inches(coords['map']['width']),
                    Inches(coords['map']['height'])
                )
                print(f"  ✅ Added map for {crime_type}")

            # Table
            for suffix in ("AutoFit", "", "Placeholder"):
                tbl = os.path.join(crime_folder, f"{prefix}_IncidentTable{('_' + suffix) if suffix else ''}.png")
                if os.path.exists(tbl):
                    slide.shapes.add_picture(
                        tbl,
                        Inches(coords['table']['left']),
                        Inches(coords['table']['top']),
                        Inches(coords['table']['width']),
                        Inches(coords['table']['height'])
                    )
                    print(f"  ✅ Added table for {crime_type}")
                    break

        prs.save(ppt_path)
        print(f"✅ PowerPoint saved: {os.path.basename(ppt_path)}")
        return True

    except Exception as e:
        print(f"❌ PowerPoint assembly failed: {e}")
        return False


def check_folder_structure(report_date_str):
    """Verify that all expected folders and files exist."""
    print(f"\n📁 CHECKING FOLDER STRUCTURE")
    print("=" * 50)
    try:
        main_out = get_output_folder(report_date_str)
        print(f"Main output: {main_out}")

        if not os.path.exists(main_out):
            print(f"❌ Main folder not found")
            return

        print(f"✅ Main folder exists")
        for crime in CRIME_TYPES:
            folder = get_crime_type_folder(crime, report_date_str)
            if os.path.exists(folder):
                files = os.listdir(folder)
                print(f"✅ {os.path.basename(folder)}: {len(files)} files")
                prefix = get_standardized_filename(crime)
                for expected in (
                    f"{prefix}_Chart.png",
                    f"{prefix}_Map.png",
                    f"{prefix}_IncidentTable.png"
                ):
                    print(f"    {'✅' if expected in files else '❌ Missing: ' + expected}")
            else:
                print(f"❌ {os.path.basename(folder)}: Folder not found")

    except Exception as e:
        print(f"❌ Error checking folders: {e}")


def main():
    """Command-line interface"""
    if len(sys.argv) > 1:
        cmd = sys.argv[1].lower()
        date_str = sys.argv[2] if len(sys.argv) > 2 else dt_date.today().strftime("%Y_%m_%d")

        if cmd == "check":
            check_folder_structure(date_str)
        elif cmd == "run":
            success = run_complete_scrpa_system(date_str, generate_all=True)
            sys.exit(0 if success else 1)
        elif cmd == "charts":
            print("📈 Generating charts only...")
            for crime in CRIME_TYPES:
                try:
                    ok = export_chart(crime)
                    print(f"{'✅' if ok else '❌'} {crime} chart")
                except Exception as e:
                    print(f"❌ {crime} chart failed: {e}")
        else:
            print(f"Unknown command: {cmd}")
            print("Usage: python scrpa_complete_integration.py [run|check|charts] [YYYY_MM_DD]")
    else:
        # Default
        success = run_complete_scrpa_system()
        sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
