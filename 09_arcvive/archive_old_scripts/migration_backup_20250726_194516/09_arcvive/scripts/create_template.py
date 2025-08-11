# create_template.py
# Save this as a file and run it

from pptx import Presentation
import os

print("Creating SCRPA PowerPoint template...")

# Create basic presentation
prs = Presentation()

# Add a title slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "SCRPA 7-Day Crime Analysis"
subtitle.text = "Hackensack Police Department"

# Add slides for each crime type
crime_types = ["MV Theft", "Burglary Auto", "Burglary - Comm & Res", "Robbery", "Sexual Offenses"]

for crime_type in crime_types:
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    from pptx.util import Inches
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = f"SCRPA Analysis - {crime_type}"

# Save template
template_dir = r"C:\Users\carucci_r\OneDrive - City of Hackensack\_25_SCRPA\TIME_Based\7_Day_Templet_SCRPA_Time"
os.makedirs(template_dir, exist_ok=True)
template_path = os.path.join(template_dir, "SCRPA_REPORT.pptx")

prs.save(template_path)
print(f"✅ Template created: {template_path}")
print("✅ Ready to run SCRPA!")