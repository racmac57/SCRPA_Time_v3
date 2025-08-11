#!/usr/bin/env python3
# 🕒 2025-07-01-17-00-00
# SCRPA_Complete_System/scrpa_complete_system.py
# Author: R. A. Carucci
# Purpose: Complete SCRPA system with charts, maps, statistics, and proper folder structure

import os
import sys
import logging
import time
import glob
import shutil
import traceback
from datetime import datetime as dt_datetime, date as dt_date
import arcpy

# --- CORE CONFIG IMPORT ---
try:
    from config import (
        get_crime_type_folder,
        get_7day_period_dates,
        CRIME_FOLDER_MAPPING,
        POWERPOINT_TEMPLATES,
        CRIME_TYPES,
        get_correct_folder_name,
        get_standardized_filename,
        get_output_folder
    )
except ImportError as e:
    print(f"❌ Config import failed: {e}")
    sys.exit(1)

# --- MAP_EXPORT DYNAMIC LOADER & STUBS ---
try:
    import map_export

    if hasattr(map_export, 'export_maps') and not hasattr(map_export, 'export_maps_for_crime_type'):
        map_export.export_maps_for_crime_type = map_export.export_maps

    if hasattr(map_export, 'export_maps'):
        export_maps = map_export.export_maps
    elif hasattr(map_export, 'export_maps_for_crime_type'):
        export_maps = map_export.export_maps_for_crime_type
    else:
        raise ImportError("No map-export function found in map_export module")

    def __stub_filter(*args, **kwargs):
        return "1=0"
    for fn in (
        'build_sql_filter_7day_excel',
        'build_sql_filter_28day_excel',
        'build_sql_filter_ytd_excel'
    ):
        if not hasattr(map_export, fn):
            setattr(map_export, fn, __stub_filter)

except ImportError as e:
    print(f"❌ map_export import failed: {e}")
    sys.exit(1)

# --- OTHER COMPONENTS ---
try:
    from chart_export import export_chart
    from incident_table_automation import (
        export_incident_table_data_with_autofit,
        create_incident_table_placeholder
    )
    from scrpa_statistical_fixed import generate_statistical_exports_fixed
    print("✅ All modules imported successfully")
except ImportError as e:
    print(f"❌ Module import failed: {e}")
    sys.exit(1)


def run_complete_scrpa_system(report_date_str=None, generate_all=True):
    """Run the full SCRPA pipeline: maps, charts, tables, stats, PPT assembly."""
    start = time.time()
    if report_date_str is None:
        report_date_str = dt_date.today().strftime("%Y_%m_%d")

    print(f"\n🚀 SCRPA COMPLETE SYSTEM\n{'='*60}")
    print(f"Report Date: {report_date_str}\nGenerate All Components: {generate_all}")

    try:
        report_date = dt_datetime.strptime(report_date_str, "%Y_%m_%d").date()
        main_out = get_output_folder(report_date_str)
        os.makedirs(main_out, exist_ok=True)

        template_ok = setup_powerpoint_template(main_out, report_date_str)

        results = {}
        total = len(CRIME_TYPES)
        print(f"\n📊 PROCESSING {total} CRIME TYPES\n{'='*50}")

        for idx, crime in enumerate(CRIME_TYPES, 1):
            print(f"\n[{idx}/{total}] Processing {crime}...")
            folder = get_crime_type_folder(crime, report_date_str)
            os.makedirs(folder, exist_ok=True)

            res = {'folder_created': True, 'map': False, 'chart': False, 'table': False}

            if generate_all:
                # Map
                print("  🗺️ Generating map...")
                try:
                    ok = export_maps(crime, folder, sys.argv)
                    res['map'] = ok
                    print(f"    {'✅' if ok else '❌'} Map: {'Success' if ok else 'Failed'}")
                except Exception:
                    print("    ❌ Map export encountered an error:")
                    traceback.print_exc()
                    res['map'] = False

                # Chart
                print("  📈 Generating chart...")
                try:
                    ok = export_chart(crime)
                    res['chart'] = ok
                    print(f"    {'✅' if ok else '❌'} Chart: {'Success' if ok else 'Failed'}")
                except Exception as e:
                    print(f"    ❌ Chart failed: {e}")
                    res['chart'] = False

                # Table
                print("  📋 Generating incident table...")
                try:
                    ok = export_incident_table_data_with_autofit(crime, report_date)
                    if not ok:
                        prefix = get_standardized_filename(crime)
                        ph = os.path.join(folder, f"{prefix}_IncidentTable_Placeholder.png")
                        ok = create_incident_table_placeholder(crime, ph)
                    res['table'] = ok
                    print(f"    {'✅' if ok else '❌'} Table: {'Success' if ok else 'Failed'}")
                except Exception as e:
                    print(f"    ❌ Table failed: {e}")
                    res['table'] = False

            results[crime] = res
            cnt = sum(1 for v in res.values() if v)
            print(f"  📊 {crime}: {cnt}/4 components successful")

        # Statistical exports
        if generate_all:
            print(f"\n📊 GENERATING STATISTICAL REPORTS\n{'-'*40}")
            try:
                ok = generate_statistical_exports_fixed(report_date)
                print(f"{'✅' if ok else '❌'} Statistical reports: {'Success' if ok else 'Failed'}")
                stats_ok = ok
            except Exception as e:
                print(f"❌ Statistical exports failed: {e}")
                stats_ok = False
        else:
            stats_ok = True

        # PPT assembly
        if generate_all and template_ok:
            print(f"\n🎯 ASSEMBLING POWERPOINT PRESENTATION\n{'-'*40}")
            try:
                ppt_ok = assemble_powerpoint_presentation(main_out, results, report_date_str)
                print(f"{'✅' if ppt_ok else '❌'} PowerPoint assembly: {'Success' if ppt_ok else 'Failed'}")
            except Exception as e:
                print(f"❌ PPT assembly failed: {e}")
                ppt_ok = False
        else:
            ppt_ok = template_ok

        # Summary
        passed = sum(1 for r in results.values() if sum(r.values()) >= 2)
        overall = (passed >= len(CRIME_TYPES)*0.6 and stats_ok and ppt_ok)
        elapsed = time.time() - start

        print(f"\n🎉 SCRPA SYSTEM COMPLETE\n{'='*60}")
        print(f"Processing time: {elapsed:.1f}s")
        print(f"Overall success: {'✅ SUCCESS' if overall else '❌ PARTIAL'}")
        print(f"Crime types processed: {passed}/{len(CRIME_TYPES)}")
        print(f"Statistical reports: {'✅' if stats_ok else '❌'}")
        print(f"PPT assembly: {'✅' if ppt_ok else '❌'}")
        print(f"Output folder: {main_out}\n\n📋 DETAILED RESULTS:")
        for crime, r in results.items():
            comps = [k for k, v in r.items() if v]
            print(f"  {crime}: {len(comps)}/4 - {', '.join(comps)}")

        return overall

    except Exception as e:
        print(f"❌ SCRPA failed: {e}")
        logging.error("SCRPA failure", exc_info=True)
        return False


def setup_powerpoint_template(out_folder, date_str):
    """Copy PPT template into output folder."""
    try:
        tpl = POWERPOINT_TEMPLATES.get("7day") or POWERPOINT_TEMPLATES.get("main")
        if not tpl or not os.path.exists(tpl):
            print(f"⚠️ PPT template not found: {tpl}")
            return False
        dest = os.path.join(out_folder, f"{get_correct_folder_name(date_str)}.pptx")
        shutil.copy2(tpl, dest)
        print(f"✅ PPT template copied: {os.path.basename(dest)}")
        return True
    except Exception as e:
        print(f"❌ PPT template setup failed: {e}")
        return False


def assemble_powerpoint_presentation(out_folder, results, date_str):
    """Add charts, maps, tables to each slide."""
    try:
        from pptx import Presentation
        from pptx.util import Inches

        pps = glob.glob(os.path.join(out_folder, "*.pptx"))
        if not pps:
            print(f"❌ No PPT found in: {out_folder}")
            return False

        prs = Presentation(pps[0])
        print(f"📝 Assembling {len(prs.slides)} slides...")

        coords = {
            'chart': {'left': 0.5, 'top': 2.8, 'width': 5.84, 'height': 4.00},
            'map':   {'left': 6.8, 'top': 2.8, 'width': 5.60, 'height': 4.03},
            'table': {'left': 6.8, 'top': 0.8, 'width': 8.33, 'height': 1.98}
        }

        for idx, crime in enumerate(CRIME_TYPES):
            if idx >= len(prs.slides):
                print(f"⚠️ No slide for {crime}")
                continue
            slide = prs.slides[idx]
            folder = get_crime_type_folder(crime, date_str)
            prefix = get_standardized_filename(crime)

            pic = os.path.join(folder, f"{prefix}_Chart.png")
            if os.path.exists(pic):
                slide.shapes.add_picture(
                    pic, Inches(coords['chart']['left']), Inches(coords['chart']['top']),
                    Inches(coords['chart']['width']), Inches(coords['chart']['height'])
                )
                print(f"  ✅ Added chart for {crime}")

            pic = os.path.join(folder, f"{prefix}_Map.png")
            if os.path.exists(pic):
                slide.shapes.add_picture(
                    pic, Inches(coords['map']['left']), Inches(coords['map']['top']),
                    Inches(coords['map']['width']), Inches(coords['map']['height'])
                )
                print(f"  ✅ Added map for {crime}")

            for suf in ("AutoFit", "", "Placeholder"):
                tbl = os.path.join(folder, f"{prefix}_IncidentTable{('_'+suf) if suf else ''}.png")
                if os.path.exists(tbl):
                    slide.shapes.add_picture(
                        tbl, Inches(coords['table']['left']), Inches(coords['table']['top']),
                        Inches(coords['table']['width']), Inches(coords['table']['height'])
                    )
                    print(f"  ✅ Added table for {crime}")
                    break

        prs.save(pps[0])
        print(f"✅ PPT saved: {os.path.basename(pps[0])}")
        return True
    except Exception as e:
        print(f"❌ PPT assembly failed: {e}")
        return False


def check_folder_structure(report_date_str):
    """Verify that all expected folders and files exist."""
    print(f"\n📁 CHECKING FOLDER STRUCTURE\n{'='*50}")
    try:
        main_out = get_output_folder(report_date_str)
        print(f"Main output: {main_out}")

        if not os.path.exists(main_out):
            print(f"❌ Main folder not found")
            return

        print(f"✅ Main folder exists")
        for crime in CRIME_TYPES:
            folder = get_crime_type_folder(crime, report_date_str)
            if os.path.exists(folder):
                files = os.listdir(folder)
                print(f"✅ {os.path.basename(folder)}: {len(files)} files")
                prefix = get_standardized_filename(crime)
                for expected in (
                    f"{prefix}_Chart.png",
                    f"{prefix}_Map.png",
                    f"{prefix}_IncidentTable.png"
                ):
                    print(f"    {'✅' if expected in files else '❌ Missing: '+expected}")
            else:
                print(f"❌ {os.path.basename(folder)}: Folder not found")
    except Exception as e:
        print(f"❌ Error checking folders: {e}")


def main():
    """CLI entrypoint."""
    if len(sys.argv) > 1:
        cmd = sys.argv[1].lower()
        date_str = sys.argv[2] if len(sys.argv) > 2 else dt_date.today().strftime("%Y_%m_%d")

        if cmd == "check":
            check_folder_structure(date_str)
        elif cmd == "run":
            ok = run_complete_scrpa_system(date_str, generate_all=True)
            sys.exit(0 if ok else 1)
        elif cmd == "charts":
            print("📈 Charts only…")
            for crime in CRIME_TYPES:
                try:
                    ok = export_chart(crime)
                    print(f"{'✅' if ok else '❌'} {crime} chart")
                except Exception as e:
                    print(f"❌ {crime} chart failed: {e}")
        else:
            print(f"Unknown command: {cmd}\nUsage: python scrpa_complete_system.py [run|check|charts] [YYYY_MM_DD]")
    else:
        success = run_complete_scrpa_system()
        sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
