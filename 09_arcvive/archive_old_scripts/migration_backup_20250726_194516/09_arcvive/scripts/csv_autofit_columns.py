# 🕒 2025-06-22-21-45-00
# Police_Data_Analysis/csv_autofit_columns
# Author: R. A. Carucci
# Purpose: Auto-fit column widths in CSV table images for PowerPoint auto-assembly with enhanced error handling and styled placeholders

import os
import logging
from datetime import date
from pptx import Presentation
from pptx.util import Inches

try:
    from config import CRIME_FOLDER_MAPPING, get_correct_folder_name, REPORT_BASE_DIR, POWERPOINT_TEMPLATES, get_7day_period_dates
    from incident_table_automation import create_formatted_csv_image_with_autofit, create_incident_table_placeholder
except ImportError as e:
    print(f"❌ Dependency error: {e}")
    print("Please ensure 'config.py' and 'incident_table_automation.py' are in the same directory")
    exit(1)

# Configure logging to match main.py
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def auto_assemble_powerpoint_with_autofit(target_date):
    """
    PowerPoint assembly prioritizing auto-fit images for Robbery and Sexual Offenses (CODE_20250622_005).
    Uses styled placeholders from incident_table_automation.py.
    """
    try:
        # Validate PowerPoint template paths
        template_path = POWERPOINT_TEMPLATES.get("main")
        fallback_template = POWERPOINT_TEMPLATES.get("7day")
        
        if not template_path or not os.path.exists(template_path):
            if fallback_template and os.path.exists(fallback_template):
                template_path = fallback_template
                logging.info(f"⚠️ Main template not found, using 7-day template: {template_path}")
            else:
                logging.error(f"❌ PowerPoint template not found: main={template_path}, 7day={fallback_template}")
                return False
        
        folder_name = get_correct_folder_name(target_date)
        base_folder = os.path.join(REPORT_BASE_DIR, folder_name)
        ppt_path = os.path.join(base_folder, f"{folder_name}.pptx")
        
        logging.info(f"📊 Auto-assembling PowerPoint with auto-fit tables: {ppt_path}")
        
        # Create PowerPoint from template
        prs = Presentation(template_path)
        
        # Process each crime type
        for i, (crime_type, folder_mapping) in enumerate(CRIME_FOLDER_MAPPING.items()):
            if i >= len(prs.slides):
                logging.warning(f"⚠️ No slide available for {crime_type} (slide index {i}), skipping")
                continue
                
            slide = prs.slides[i]
            crime_folder = os.path.join(base_folder, folder_mapping)
            os.makedirs(crime_folder, exist_ok=True)
            
            logging.info(f"🔄 Processing slide {i+1}: {crime_type}")
            
            # Add chart
            chart_file = os.path.join(crime_folder, f"{folder_mapping}_Chart.png")
            if os.path.exists(chart_file):
                slide.shapes.add_picture(chart_file, Inches(0.5), Inches(2), Inches(6), Inches(4))
                logging.info(f"  ✅ Inserted chart: {chart_file}")
            else:
                logging.warning(f"  ⚠️ Chart file not found: {chart_file}")
            
            # Add map
            map_file = os.path.join(crime_folder, f"{folder_mapping}_Map.png")
            if os.path.exists(map_file):
                slide.shapes.add_picture(map_file, Inches(7), Inches(2), Inches(6), Inches(4))
                logging.info(f"  ✅ Inserted map: {map_file}")
            else:
                logging.warning(f"  ⚠️ Map file not found: {map_file}")
            
            # Add incident table (auto-fit for Robbery and Sexual Offenses)
            autofit_crimes = ["Robbery", "Sexual Offenses"]
            
            if crime_type in autofit_crimes:
                autofit_image = os.path.join(crime_folder, f"{folder_mapping}_Incidents_Table_AutoFit.png")
                if os.path.exists(autofit_image):
                    slide.shapes.add_picture(autofit_image, Inches(0.5), Inches(6.5), Inches(12), Inches(2))
                    logging.info(f"  ✅ Inserted AUTO-FIT incident table: {autofit_image}")
                    continue
            
            # Fallback to CSV or placeholder
            incident_files = [
                os.path.join(crime_folder, f"{folder_mapping}_Incidents_Manual_Entry.csv"),
                os.path.join(crime_folder, f"{folder_mapping}_Incidents_Placeholder.png")
            ]
            
            for incident_file in incident_files:
                if os.path.exists(incident_file):
                    if incident_file.endswith('.csv'):
                        image_path = incident_file.replace('.csv', '_Table.png')
                        if create_formatted_csv_image_with_autofit(incident_file, image_path):
                            slide.shapes.add_picture(image_path, Inches(0.5), Inches(6.5), Inches(12), Inches(2))
                            logging.info(f"  ✅ Inserted converted CSV table: {image_path}")
                            break
                    else:
                        slide.shapes.add_picture(incident_file, Inches(0.5), Inches(6.5), Inches(12), Inches(2))
                        logging.info(f"  ✅ Inserted placeholder: {incident_file}")
                        break
            else:
                # Generate placeholder if no CSV or existing placeholder is found
                logging.warning(f"  ⚠️ No incident table or placeholder found for {crime_type}, generating new placeholder")
                start_date, end_date = get_7day_period_dates(target_date)
                placeholder_path = os.path.join(crime_folder, f"{folder_mapping}_Incidents_Placeholder.png")
                if create_incident_table_placeholder(crime_type, placeholder_path, start_date, end_date):
                    slide.shapes.add_picture(placeholder_path, Inches(0.5), Inches(6.5), Inches(12), Inches(2))
                    logging.info(f"  ✅ Generated and inserted new placeholder: {placeholder_path}")
        
        # Save PowerPoint
        os.makedirs(base_folder, exist_ok=True)
        prs.save(ppt_path)
        logging.info(f"✅ PowerPoint saved: {ppt_path}")
        return True
        
    except Exception as e:
        logging.error(f"❌ PowerPoint assembly failed: {e}")
        return False

if __name__ == "__main__":
    try:
        import pptx
    except ImportError:
        print("❌ Missing dependency: python-pptx")
        print("Please install it using: pip install python-pptx")
        exit(1)
    
    from datetime import date
    from config import CRIME_FOLDER_MAPPING, get_correct_folder_name, REPORT_BASE_DIR
    
    logging.info("🧪 Testing PowerPoint assembly")
    logging.info("=" * 50)
    
    target_date = date(2025, 6, 17)
    
    success = auto_assemble_powerpoint_with_autofit(target_date)
    
    if success:
        logging.info("✅ PowerPoint assembly test completed successfully")
        folder_name = get_correct_folder_name(target_date)
        base_folder = os.path.join(REPORT_BASE_DIR, folder_name)
        logging.info(f"\n📋 Check these files:")
        logging.info(f"PowerPoint: {base_folder}\\{folder_name}.pptx")
        
        for crime_type, folder_mapping in CRIME_FOLDER_MAPPING.items():
            crime_folder = os.path.join(base_folder, folder_mapping)
            logging.info(f"\n{crime_type} folder: {crime_folder}")
            try:
                files = os.listdir(crime_folder)
                for file in files:
                    logging.info(f"  - {file}")
            except:
                logging.info(f"  - Folder not found")
    else:
        logging.error("❌ PowerPoint assembly test failed")